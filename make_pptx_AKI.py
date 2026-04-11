from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ─────────────────────────────────────────────
# Color Palette  (Pretendard / Indigo theme)
# ─────────────────────────────────────────────
INDIGO      = RGBColor(46, 61, 134)
LIGHT_BLUE  = RGBColor(204, 210, 240)
BLACK       = RGBColor(0, 0, 0)
WHITE       = RGBColor(255, 255, 255)
RED         = RGBColor(192, 0, 0)
PALE_INDIGO = RGBColor(232, 235, 248)
FONT_MAIN   = "Pretendard"

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(blank_layout)

def add_bg(slide, color=WHITE):
    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = color; bg.line.fill.background()

def add_header(slide, chapter_num, section_title, subtitle="Medical Student Case Presentation"):
    HEADER_H = Inches(0.95)
    hbar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, HEADER_H)
    hbar.fill.solid(); hbar.fill.fore_color.rgb = WHITE; hbar.line.fill.background()

    rule = slide.shapes.add_shape(1, Inches(0.5), HEADER_H - Pt(2), prs.slide_width - Inches(1.0), Pt(2))
    rule.fill.solid(); rule.fill.fore_color.rgb = INDIGO; rule.line.fill.background()

    num_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.05), Inches(0.8), Inches(0.85))
    tf = num_box.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    run = p.add_run(); run.text = chapter_num
    run.font.name = FONT_MAIN; run.font.bold = True
    run.font.size = Pt(40); run.font.color.rgb = INDIGO

    txt_box = slide.shapes.add_textbox(Inches(1.25), Inches(0.08), Inches(11.5), Inches(0.38))
    tf2 = txt_box.text_frame; tf2.word_wrap = False
    p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.LEFT
    run2 = p2.add_run(); run2.text = subtitle
    run2.font.name = FONT_MAIN; run2.font.size = Pt(13); run2.font.color.rgb = BLACK

    sec_box = slide.shapes.add_textbox(Inches(1.25), Inches(0.46), Inches(11.5), Inches(0.42))
    tf3 = sec_box.text_frame; tf3.word_wrap = False
    p3 = tf3.paragraphs[0]; p3.alignment = PP_ALIGN.LEFT
    run3 = p3.add_run(); run3.text = section_title
    run3.font.name = FONT_MAIN; run3.font.bold = True
    run3.font.size = Pt(18); run3.font.color.rgb = INDIGO

    foot = slide.shapes.add_textbox(Inches(11.8), Inches(7.2), Inches(1.3), Inches(0.25))
    tf4 = foot.text_frame
    p4 = tf4.paragraphs[0]; p4.alignment = PP_ALIGN.RIGHT
    run4 = p4.add_run(); run4.text = chapter_num
    run4.font.name = FONT_MAIN; run4.font.size = Pt(11); run4.font.color.rgb = LIGHT_BLUE

def add_textbox(slide, text, left, top, width, height,
                font_size=16, bold=False, color=BLACK,
                align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.name = FONT_MAIN; run.font.size = Pt(font_size)
    run.font.bold = bold; run.font.color.rgb = color
    return txBox

def add_bullets(slide, items, left, top, width, height, font_size=14):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = item.get('level', 0)
        run = p.add_run()
        run.text = item['text']
        run.font.name = FONT_MAIN
        run.font.size = Pt(item.get('size', font_size))
        run.font.bold = item.get('bold', False)
        run.font.color.rgb = item.get('color', BLACK)
        p.space_after = Pt(3)

def add_table(slide, headers, rows, left, top, width, col_widths=None):
    rows_count = len(rows) + 1
    cols_count = len(headers)
    tbl = slide.shapes.add_table(
        rows_count, cols_count, left, top, width,
        Inches(0.33 * rows_count + 0.33)
    ).table
    if col_widths:
        for i, w in enumerate(col_widths): tbl.columns[i].width = Inches(w)
    for i, h in enumerate(headers):
        cell = tbl.cell(0, i); cell.fill.solid(); cell.fill.fore_color.rgb = INDIGO
        p = cell.text_frame.paragraphs[0]; run = p.add_run()
        run.text = h; run.font.name = FONT_MAIN; run.font.bold = True
        run.font.size = Pt(12); run.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            if (ri + 1) % 2 == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = PALE_INDIGO
            p = cell.text_frame.paragraphs[0]; run = p.add_run()
            is_ab = val.get('abnormal', False) if isinstance(val, dict) else False
            text  = val['text'] if isinstance(val, dict) else str(val)
            run.text = text; run.font.name = FONT_MAIN; run.font.size = Pt(12)
            run.font.color.rgb = RED if is_ab else BLACK
            run.font.bold = is_ab; p.alignment = PP_ALIGN.CENTER

def add_info_box(slide, left, top, width, height):
    box = slide.shapes.add_shape(1, left, top, width, height)
    box.fill.solid(); box.fill.fore_color.rgb = PALE_INDIGO; box.line.color.rgb = INDIGO
    return box


# ═══════════════════════════════════════════════════════════
# SLIDE 1: Title
# ═══════════════════════════════════════════════════════════
s = add_slide()
add_bg(s, INDIGO)

deco = s.shapes.add_shape(1, Inches(9.8), 0, Inches(3.53), prs.slide_height)
deco.fill.solid(); deco.fill.fore_color.rgb = RGBColor(36, 51, 124); deco.line.fill.background()

acc = s.shapes.add_shape(1, Inches(0.5), Inches(2.85), Inches(9.0), Pt(2))
acc.fill.solid(); acc.fill.fore_color.rgb = LIGHT_BLUE; acc.line.fill.background()

add_textbox(s, "A Case of Acute Kidney Injury (AKI)",
            Inches(0.5), Inches(1.3), Inches(9.0), Inches(1.2),
            font_size=30, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_textbox(s, "Urosepsis-induced AKI  |  KDIGO Stage 3",
            Inches(0.5), Inches(2.5), Inches(9.0), Inches(0.5),
            font_size=17, color=LIGHT_BLUE, align=PP_ALIGN.LEFT)
add_textbox(s, "Department of Nephrology / Internal Medicine",
            Inches(0.5), Inches(3.15), Inches(9.0), Inches(0.45),
            font_size=14, color=WHITE, align=PP_ALIGN.LEFT)
add_textbox(s, "Medical Student Case Presentation  |  April 2026",
            Inches(0.5), Inches(3.65), Inches(9.0), Inches(0.4),
            font_size=13, color=LIGHT_BLUE, align=PP_ALIGN.LEFT)
add_textbox(s, "CASE\nPRESENTATION",
            Inches(10.1), Inches(2.8), Inches(3.0), Inches(2.0),
            font_size=28, bold=True, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════
# SLIDE 2: Patient Information
# ═══════════════════════════════════════════════════════════
s = add_slide()
add_bg(s)
add_header(s, "01", "Patient Information")

add_table(s,
    ["Field", "Details"],
    [
        ["Age / Sex", "75 y/o Female"],
        ["Chief Complaint", "Decreased urine output, fever, altered mental status — 2 days"],
        ["Date of Visit", "Emergency admission via 119 ambulance"],
        ["Triage Level", {"text": "ESI Level 2 (Emergent)", "abnormal": True}],
        ["Vital Signs on Arrival",
         {"text": "BP 88/56  |  HR 118  |  RR 24  |  Temp 38.9°C  |  SpO₂ 96% (RA)", "abnormal": True}],
    ],
    Inches(0.5), Inches(1.1), Inches(12.3),
    col_widths=[3.2, 9.1]
)

# ═══════════════════════════════════════════════════════════
# SLIDE 3: HPI
# ═══════════════════════════════════════════════════════════
s = add_slide()
add_bg(s)
add_header(s, "02", "History of Present Illness (HPI)")

add_bullets(s, [
    {'text': 'Timeline', 'bold': True, 'size': 15, 'color': INDIGO},
    {'text': '5 days prior — dysuria, urinary frequency, mild flank pain (no treatment sought)', 'level': 1, 'size': 13},
    {'text': '3 days prior — high fever (38.9°C), chills, nausea', 'level': 1, 'size': 13},
    {'text': '2 days prior — oliguria onset (<400 mL/day), progressive confusion, family called 119', 'level': 1, 'size': 13},
    {'text': '', 'size': 6},
    {'text': 'OLDCARTS', 'bold': True, 'size': 15, 'color': INDIGO},
], Inches(0.5), Inches(1.1), Inches(5.5), Inches(5.8))

add_table(s,
    ["Feature", "Description"],
    [
        ["Onset",         "Gradual, 5 days ago"],
        ["Location",      "Bilateral flank / suprapubic pain"],
        ["Duration",      {"text": "5 days, worsening", "abnormal": True}],
        ["Character",     "Dull flank ache → oliguria → AMS"],
        ["Radiation",     "None"],
        ["Severity",      {"text": "Oliguria 200 mL/24h", "abnormal": True}],
        ["Associated Sx", {"text": "Fever, chills, nausea, confusion, leg edema", "abnormal": True}],
    ],
    Inches(6.0), Inches(1.1), Inches(7.0),
    col_widths=[2.2, 4.8]
)

# ═══════════════════════════════════════════════════════════
# SLIDE 4: PMH / FH / SH
# ═══════════════════════════════════════════════════════════
s = add_slide()
add_bg(s)
add_header(s, "03", "Past Medical / Family / Social History")

add_bullets(s, [
    {'text': 'Past Medical History', 'bold': True, 'size': 15, 'color': INDIGO},
    {'text': 'T2DM (15 yrs) — Metformin 1000mg BID + Glimepiride 2mg QD', 'level': 1, 'size': 13},
    {'text': 'Hypertension (12 yrs) — Losartan 50mg QD (ARB)', 'level': 1, 'size': 13},
    {'text': 'CKD Stage 2  (baseline Cr 1.1 mg/dL, 6 months ago)', 'level': 1, 'size': 13, 'color': RED},
    {'text': 'Recurrent UTI — 2 episodes in past year', 'level': 1, 'size': 13},
    {'text': 'No prior surgeries  |  NKDA', 'level': 1, 'size': 13},
    {'text': '', 'size': 5},
    {'text': 'Family History', 'bold': True, 'size': 15, 'color': INDIGO},
    {'text': 'Mother: CKD (ESRD, HD-dependent) — deceased at 78', 'level': 1, 'size': 13, 'color': RED},
    {'text': 'Daughter: T2DM', 'level': 1, 'size': 13},
    {'text': '', 'size': 5},
    {'text': 'Social History', 'bold': True, 'size': 15, 'color': INDIGO},
    {'text': 'Non-smoker, non-drinker', 'level': 1, 'size': 13},
    {'text': 'Lives alone; poor oral intake × 3 days (poor PO hydration)', 'level': 1, 'size': 13, 'color': RED},
    {'text': 'No NSAID / contrast exposure in past 2 weeks', 'level': 1, 'size': 13},
], Inches(0.5), Inches(1.1), Inches(12.5), Inches(6.0))

# ═══════════════════════════════════════════════════════════
# SLIDE 5: ROS
# ═══════════════════════════════════════════════════════════
s = add_slide()
add_bg(s)
add_header(s, "04", "Review of Systems (ROS)")

add_table(s,
    ["System", "Positive (+)", "Negative (−)"],
    [
        ["Constitutional", {"text": "Fever 38.9°C, chills, fatigue, poor PO intake", "abnormal": True}, "No weight loss"],
        ["HEENT",          "—", "Unremarkable"],
        ["Cardiovascular", {"text": "Peripheral edema (bilateral lower extremity)", "abnormal": True}, "No chest pain, no palpitations"],
        ["Pulmonary",      {"text": "Mild dyspnea", "abnormal": True}, "No cough, no hemoptysis"],
        ["GI",             {"text": "Nausea, decreased appetite", "abnormal": True}, "No vomiting, no diarrhea, no melena"],
        ["GU",             {"text": "Dysuria, frequency (5d ago); oliguria 200mL/24h", "abnormal": True}, "No hematuria (gross)"],
        ["MSK",            "—", "No joint pain"],
        ["Neuro",          {"text": "Confusion, lethargy (AMS)", "abnormal": True}, "No focal neurologic deficit"],
        ["Skin",           "—", "No rash, no petechiae"],
    ],
    Inches(0.5), Inches(1.1), Inches(12.3),
    col_widths=[2.1, 5.0, 5.2]
)

# ═══════════════════════════════════════════════════════════
# SLIDE 6: Physical Exam
# ═══════════════════════════════════════════════════════════
s = add_slide()
add_bg(s)
add_header(s, "05", "Physical Examination")

add_table(s,
    ["BP", "HR", "RR", "Temp", "SpO₂"],
    [[
        {"text": "88/56 mmHg", "abnormal": True},
        {"text": "118 bpm",    "abnormal": True},
        {"text": "24 /min",    "abnormal": True},
        {"text": "38.9 °C",   "abnormal": True},
        "96% (RA)",
    ]],
    Inches(0.5), Inches(1.1), Inches(12.3),
    col_widths=[2.8, 2.2, 2.2, 2.2, 2.9]
)

add_table(s,
    ["System", "Findings"],
    [
        ["General",       {"text": "Ill-appearing, lethargic, diaphoretic; GCS 13 (E3V4M6)", "abnormal": True}],
        ["HEENT / Neck",  "Dry mucous membranes, no JVD, no lymphadenopathy"],
        ["Chest / Lungs", "Clear to auscultation bilaterally; no crackles"],
        ["Heart",         {"text": "Tachycardic, regular rhythm; no murmur", "abnormal": True}],
        ["Abdomen",       {"text": "Suprapubic tenderness (+), CVA tenderness bilateral (+)", "abnormal": True}],
        ["Extremities",   {"text": "Bilateral pitting edema 2+ (up to knee)", "abnormal": True}],
        ["Neuro",         {"text": "Confusion, disoriented to time & place; no focal deficits", "abnormal": True}],
    ],
    Inches(0.5), Inches(2.38), Inches(12.3),
    col_widths=[2.5, 9.8]
)

# ═══════════════════════════════════════════════════════════
# SLIDE 7: Labs & Imaging
# ═══════════════════════════════════════════════════════════
s = add_slide()
add_bg(s)
add_header(s, "06", "Laboratory & Imaging Results")

add_table(s,
    ["Lab", "Result", "Ref", ""],
    [
        ["Creatinine",  {"text": "4.8 mg/dL",   "abnormal": True},  "0.5–1.1",  {"text": "↑↑↑", "abnormal": True}],
        ["BUN",         {"text": "68 mg/dL",     "abnormal": True},  "7–20",     {"text": "↑↑",  "abnormal": True}],
        ["K⁺",          {"text": "5.8 mEq/L",   "abnormal": True},  "3.5–5.0",  {"text": "↑",   "abnormal": True}],
        ["Bicarbonate", {"text": "14 mEq/L",     "abnormal": True},  "22–29",    {"text": "↓",   "abnormal": True}],
        ["Na⁺",         "132 mEq/L",   "135–145", "↓"],
        ["WBC",         {"text": "18.4 ×10³/μL", "abnormal": True},  "4–10",     {"text": "↑↑",  "abnormal": True}],
        ["CRP",         {"text": "186 mg/L",      "abnormal": True},  "<5",       {"text": "↑↑↑", "abnormal": True}],
        ["Lactate",     {"text": "3.8 mmol/L",   "abnormal": True},  "<2.0",     {"text": "↑↑",  "abnormal": True}],
        ["Urine output","200 mL/24h",  ">400 mL", "Oliguria"],
    ],
    Inches(0.5), Inches(1.1), Inches(6.3),
    col_widths=[2.1, 1.9, 1.3, 1.0]
)

add_bullets(s, [
    {'text': 'Urinalysis', 'bold': True, 'size': 14, 'color': INDIGO},
    {'text': 'Pyuria (WBC >100/hpf), bacteriuria', 'level': 1, 'size': 13, 'color': RED, 'bold': True},
    {'text': 'Muddy brown granular casts (ATN)', 'level': 1, 'size': 13, 'color': RED, 'bold': True},
    {'text': 'Urine Na: 68 mEq/L  |  FENa: 3.2%  (>2% → intrinsic)', 'level': 1, 'size': 13},
    {'text': '', 'size': 5},
    {'text': 'Blood / Urine Culture', 'bold': True, 'size': 14, 'color': INDIGO},
    {'text': 'E. coli — blood culture (+) × 2 sets', 'level': 1, 'size': 13, 'color': RED, 'bold': True},
    {'text': 'E. coli — urine culture (+), >10⁵ CFU/mL', 'level': 1, 'size': 13, 'color': RED},
    {'text': '', 'size': 5},
    {'text': 'Renal Ultrasound', 'bold': True, 'size': 14, 'color': INDIGO},
    {'text': 'Bilateral kidney size normal (R: 10.8cm, L: 11.0cm)', 'level': 1, 'size': 13},
    {'text': 'No hydronephrosis, no obstruction', 'level': 1, 'size': 13},
    {'text': 'Increased echogenicity bilaterally (parenchymal edema)', 'level': 1, 'size': 13, 'color': RED},
], Inches(6.6), Inches(1.1), Inches(6.4), Inches(5.8))

# ═══════════════════════════════════════════════════════════
# SLIDE 8: Problem List & Assessment
# ═══════════════════════════════════════════════════════════
s = add_slide()
add_bg(s)
add_header(s, "07", "Problem List & Assessment")

add_bullets(s, [
    {'text': 'Problem List', 'bold': True, 'size': 16, 'color': INDIGO},
    {'text': '1.  AKI Stage 3 (KDIGO) — Cr 4.8 from baseline 1.1 (PRIMARY)', 'level': 1, 'color': RED, 'bold': True, 'size': 14},
    {'text': '2.  Urosepsis — E. coli bacteremia (septic shock: MAP <65)',     'level': 1, 'color': RED, 'size': 14},
    {'text': '3.  Hyperkalemia — K⁺ 5.8 mEq/L',                               'level': 1, 'color': RED, 'size': 14},
    {'text': '4.  Metabolic acidosis — HCO₃ 14 mEq/L',                        'level': 1, 'size': 14},
    {'text': '5.  Underlying CKD Stage 2  (acute-on-chronic)',                 'level': 1, 'size': 14},
    {'text': '6.  T2DM / HTN — medications to hold (Metformin, ARB)',          'level': 1, 'size': 14},
], Inches(0.5), Inches(1.1), Inches(12.3), Inches(3.0))

box = add_info_box(s, Inches(0.5), Inches(4.45), Inches(12.3), Inches(2.6))
tf = box.text_frame; tf.word_wrap = True
tf.margin_left = Inches(0.18); tf.margin_top = Inches(0.12)
p = tf.paragraphs[0]
run = p.add_run(); run.text = "Assessment"
run.font.name = FONT_MAIN; run.font.bold = True
run.font.size = Pt(14); run.font.color.rgb = INDIGO
p2 = tf.add_paragraph()
run2 = p2.add_run()
run2.text = (
    "75 y/o female with CKD stage 2, T2DM, and HTN on ARB, presenting with 5-day history of UTI "
    "progressing to urosepsis (E. coli bacteremia) with septic shock (BP 88/56, lactate 3.8) and "
    "oliguric AKI KDIGO Stage 3 (Cr 4.8 from baseline 1.1). Urinalysis showing muddy brown casts "
    "and elevated FENa (3.2%) confirm intrinsic AKI (sepsis-induced ATN). Urgent nephrology consult "
    "and RRT evaluation required."
)
run2.font.name = FONT_MAIN; run2.font.size = Pt(12); run2.font.color.rgb = BLACK

# ═══════════════════════════════════════════════════════════
# SLIDE 9: Differential Diagnosis
# ═══════════════════════════════════════════════════════════
s = add_slide()
add_bg(s)
add_header(s, "08", "Differential Diagnosis")

add_table(s,
    ["Diagnosis", "Supporting Evidence", "Against"],
    [
        [{"text": "Sepsis-induced ATN ✓", "abnormal": True},
         "E. coli bacteremia, septic shock, muddy brown casts, FENa >2%",
         "—"],
        ["Pre-renal AKI",
         "Hypotension, poor PO intake, dehydration",
         "FENa 3.2% (>2%), muddy brown casts → intrinsic, not pre-renal"],
        ["Obstructive (Post-renal) AKI",
         "Bilateral flank pain, older female",
         "No hydronephrosis on U/S, no stone, no mass"],
        ["Drug-induced AKI (ARB)",
         "On Losartan, pre-existing CKD",
         "No NSAID/contrast, primary driver is sepsis not drug"],
        ["Rapidly Progressive GN",
         "Acute Cr rise, CKD background",
         "No hematuria, no red cell casts, infectious source identified"],
    ],
    Inches(0.5), Inches(1.1), Inches(12.3),
    col_widths=[2.9, 4.5, 4.9]
)

# ═══════════════════════════════════════════════════════════
# SLIDE 10: Final Diagnosis
# ═══════════════════════════════════════════════════════════
s = add_slide()
add_bg(s)
add_header(s, "09", "Final Diagnosis")

add_bullets(s, [
    {'text': 'Acute Kidney Injury (AKI) — KDIGO Stage 3',
     'bold': True, 'size': 22, 'color': RED},
    {'text': 'Etiology: Sepsis-induced Acute Tubular Necrosis (ATN)  |  Source: Urosepsis (E. coli)',
     'bold': True, 'size': 15, 'color': INDIGO},
    {'text': '', 'size': 7},
    {'text': 'KDIGO Staging Criteria (2012)', 'bold': True, 'size': 15, 'color': INDIGO},
    {'text': '✅  Cr rise ≥ 3× baseline  (1.1 → 4.8 mg/dL = 4.4× within 48h)', 'level': 1, 'size': 14},
    {'text': '✅  Oliguria < 0.3 mL/kg/hr × ≥ 24h  (200 mL/24h in 60 kg patient = 0.14 mL/kg/hr)', 'level': 1, 'size': 14},
    {'text': '  → KDIGO Stage 3 (highest severity)', 'level': 1, 'size': 14, 'color': RED, 'bold': True},
    {'text': '', 'size': 7},
    {'text': 'Sepsis Criteria (Sepsis-3)', 'bold': True, 'size': 15, 'color': INDIGO},
    {'text': 'SOFA score ≥ 2  |  Suspected infection (UTI → bacteremia) → Sepsis', 'level': 1, 'size': 14},
    {'text': 'MAP < 65 mmHg requiring vasopressor + lactate > 2 mmol/L → Septic Shock', 'level': 1, 'size': 14, 'color': RED},
    {'text': '', 'size': 7},
    {'text': 'Acute-on-Chronic Kidney Disease (AoC-CKD)  on background CKD Stage 2', 'level': 0, 'size': 13, 'color': INDIGO},
], Inches(0.5), Inches(1.1), Inches(12.3), Inches(6.0))

# ═══════════════════════════════════════════════════════════
# SLIDE 11: Treatment Plan
# ═══════════════════════════════════════════════════════════
s = add_slide()
add_bg(s)
add_header(s, "10", "Treatment Plan")

add_bullets(s, [
    {'text': 'Immediate / ICU-level Management', 'bold': True, 'size': 14, 'color': RED},
    {'text': 'Aggressive IV fluid resuscitation (NS or balanced crystalloid 30 mL/kg bolus) → target MAP ≥ 65 mmHg',
     'level': 1, 'size': 12, 'bold': True},
    {'text': 'Vasopressor: Norepinephrine if fluid-refractory  |  Hourly urine output monitoring (foley catheter)',
     'level': 1, 'size': 12},
    {'text': 'HOLD: Metformin (lactic acidosis risk), Losartan (reduces GFR)', 'level': 1, 'size': 12, 'color': RED},
], Inches(0.5), Inches(1.1), Inches(12.3), Inches(1.6))

add_table(s,
    ["Drug / Intervention", "Dose / Regimen", "Route", "Indication"],
    [
        [{"text": "Ceftriaxone", "abnormal": False}, "2g q24h (renally adjusted)", "IV", "E. coli urosepsis"],
        ["Norepinephrine",   "0.01–0.3 mcg/kg/min", "IV", "Septic shock (MAP<65)"],
        [{"text": "Calcium gluconate", "abnormal": False}, "1g over 10 min  (STAT)", "IV", {"text": "Hyperkalemia K⁺5.8 — cardiac protection", "abnormal": True}],
        ["Sodium bicarbonate", "1–2 mEq/kg",         "IV", "Metabolic acidosis + K⁺ redistribution"],
        ["Insulin + D50W",   "10U Regular + 50mL",   "IV", "Hyperkalemia — K⁺ shift"],
        ["Kayexalate (SPS)", "15–30g QD/BID",        "PO/PR", "K⁺ elimination"],
    ],
    Inches(0.5), Inches(2.8), Inches(12.3),
    col_widths=[2.9, 3.3, 1.2, 4.9]
)

add_bullets(s, [
    {'text': 'RRT Indications (consider urgent CRRT/HD)', 'bold': True, 'size': 13, 'color': RED},
    {'text': 'Refractory hyperkalemia  |  Severe metabolic acidosis  |  Oliguria / fluid overload  |  Uremic symptoms',
     'level': 1, 'size': 12},
    {'text': 'Follow-up: Nephrology consult (STAT)  |  Serial Cr/K⁺ q6–8h  |  ID consult  |  Repeat U/S in 24h',
     'level': 1, 'size': 12},
], Inches(0.5), Inches(6.15), Inches(12.3), Inches(1.2))

# ═══════════════════════════════════════════════════════════
# SLIDE 12: Discussion & Key Takeaways
# ═══════════════════════════════════════════════════════════
s = add_slide()
add_bg(s)
add_header(s, "11", "Discussion & Key Takeaways")

add_bullets(s, [
    {'text': 'KDIGO AKI Staging: Diagnose Early, Intervene Fast', 'bold': True, 'size': 16, 'color': RED},
    {'text': 'Stage 1: Cr ×1.5 or +0.3  |  Stage 2: Cr ×2  |  Stage 3: Cr ×3 or ≥4.0 or RRT',
     'level': 1, 'size': 12},
    {'text': '', 'size': 5},
    {'text': 'FENa > 2% = Intrinsic (ATN)  |  FENa < 1% = Pre-renal', 'bold': True, 'size': 16, 'color': INDIGO},
    {'text': 'FENa = (Urine Na × Serum Cr) / (Serum Na × Urine Cr) × 100%', 'level': 1, 'size': 12},
    {'text': 'Caveat: FENa unreliable with diuretics → use FEUrea instead', 'level': 1, 'size': 12},
    {'text': '', 'size': 5},
    {'text': 'Hyperkalemia in AKI: Treat Empirically — K⁺ > 5.5 is an emergency', 'bold': True, 'size': 16, 'color': INDIGO},
    {'text': 'Step 1: Membrane stabilization (Ca gluconate)  |  Step 2: Shift (Insulin+D50, NaHCO₃)',
     'level': 1, 'size': 12},
    {'text': 'Step 3: Elimination (Kayexalate, dialysis) — target K⁺ < 5.0', 'level': 1, 'size': 12},
    {'text': '', 'size': 5},
    {'text': 'Nephrotoxins to Avoid / Stop in AKI', 'bold': True, 'size': 16, 'color': INDIGO},
    {'text': 'HOLD: ACEi/ARB, NSAIDs, Metformin, aminoglycosides, IV contrast (unless life-saving)',
     'level': 1, 'size': 12},
    {'text': '', 'size': 5},
    {'text': 'Sepsis-induced AKI Recovery: Often reversible if source controlled early', 'bold': True, 'size': 15, 'color': INDIGO},
    {'text': 'Cr nadir usually 7–14 days post-sepsis; 30–50% progress to CKD advancement if delayed Tx',
     'level': 1, 'size': 12},
], Inches(0.5), Inches(1.1), Inches(12.3), Inches(6.0))

# ─────────────────────────────────────────────
# Save
# ─────────────────────────────────────────────
output_path = "/home/user/endo/case_AKI.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
