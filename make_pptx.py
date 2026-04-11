from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ─────────────────────────────────────────────
# Color Palette (Pretendard / Indigo theme — based on reference design)
# ─────────────────────────────────────────────
INDIGO      = RGBColor(46, 61, 134)     # #2E3D86  — primary brand color
LIGHT_BLUE  = RGBColor(204, 210, 240)   # #CCD2F0  — soft accent / subtext
BLACK       = RGBColor(0, 0, 0)         # #000000  — body text
WHITE       = RGBColor(255, 255, 255)   # #FFFFFF
RED         = RGBColor(192, 0, 0)       # abnormal lab values
LIGHT_BG    = RGBColor(248, 249, 252)   # near-white background
PALE_INDIGO = RGBColor(232, 235, 248)   # table stripe / box fill

FONT_MAIN   = "Pretendard"

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]   # blank


# ─────────────────────────────────────────────
# Helper: new slide
# ─────────────────────────────────────────────
def add_slide():
    return prs.slides.add_slide(blank_layout)


# ─────────────────────────────────────────────
# Helper: solid background rectangle
# ─────────────────────────────────────────────
def add_bg(slide, color=WHITE):
    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()


# ─────────────────────────────────────────────
# Helper: top header bar  (matches reference style)
#   - left:  large chapter number in INDIGO
#   - right of number: small subtitle + bold section title
#   - thin INDIGO rule under header
# ─────────────────────────────────────────────
def add_header(slide, chapter_num: str, section_title: str, subtitle: str = "Medical Student Case Presentation"):
    HEADER_H = Inches(0.95)

    # White header background
    hbar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, HEADER_H)
    hbar.fill.solid()
    hbar.fill.fore_color.rgb = WHITE
    hbar.line.fill.background()

    # Thin INDIGO rule at bottom of header
    rule = slide.shapes.add_shape(1, Inches(0.5), HEADER_H - Pt(2), prs.slide_width - Inches(1.0), Pt(2))
    rule.fill.solid()
    rule.fill.fore_color.rgb = INDIGO
    rule.line.fill.background()

    # Chapter number (large, INDIGO)
    num_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.05), Inches(0.8), Inches(0.85))
    tf = num_box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = chapter_num
    run.font.name = FONT_MAIN
    run.font.bold = True
    run.font.size = Pt(40)
    run.font.color.rgb = INDIGO

    # Subtitle line (smaller, black)
    txt_box = slide.shapes.add_textbox(Inches(1.25), Inches(0.08), Inches(11.5), Inches(0.38))
    tf2 = txt_box.text_frame
    tf2.word_wrap = False
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.LEFT
    run2 = p2.add_run()
    run2.text = subtitle
    run2.font.name = FONT_MAIN
    run2.font.size = Pt(13)
    run2.font.color.rgb = BLACK

    # Section title (larger, INDIGO, semibold)
    sec_box = slide.shapes.add_textbox(Inches(1.25), Inches(0.46), Inches(11.5), Inches(0.42))
    tf3 = sec_box.text_frame
    tf3.word_wrap = False
    p3 = tf3.paragraphs[0]
    p3.alignment = PP_ALIGN.LEFT
    run3 = p3.add_run()
    run3.text = section_title
    run3.font.name = FONT_MAIN
    run3.font.bold = True
    run3.font.size = Pt(18)
    run3.font.color.rgb = INDIGO

    # Footer page indicator (bottom right)
    foot = slide.shapes.add_textbox(Inches(11.8), Inches(7.2), Inches(1.3), Inches(0.25))
    tf4 = foot.text_frame
    p4 = tf4.paragraphs[0]
    p4.alignment = PP_ALIGN.RIGHT
    run4 = p4.add_run()
    run4.text = chapter_num
    run4.font.name = FONT_MAIN
    run4.font.size = Pt(11)
    run4.font.color.rgb = LIGHT_BLUE


# ─────────────────────────────────────────────
# Helper: textbox
# ─────────────────────────────────────────────
def add_textbox(slide, text, left, top, width, height,
                font_size=16, bold=False, color=BLACK,
                align=PP_ALIGN.LEFT, wrap=True, font_name=FONT_MAIN):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


# ─────────────────────────────────────────────
# Helper: bullet list
# ─────────────────────────────────────────────
def add_bullets(slide, items, left, top, width, height, font_size=14):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.level = item.get('level', 0)
        run = p.add_run()
        run.text = item['text']
        run.font.name = FONT_MAIN
        run.font.size = Pt(item.get('size', font_size))
        run.font.bold = item.get('bold', False)
        run.font.color.rgb = item.get('color', BLACK)
        p.space_after = Pt(3)


# ─────────────────────────────────────────────
# Helper: table
#   rows: list of lists; cell can be str or dict{'text','abnormal'}
# ─────────────────────────────────────────────
def add_table(slide, headers, rows, left, top, width, col_widths=None):
    rows_count = len(rows) + 1
    cols_count = len(headers)
    tbl = slide.shapes.add_table(
        rows_count, cols_count, left, top, width,
        Inches(0.33 * rows_count + 0.33)
    ).table

    if col_widths:
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = Inches(w)

    # Header row
    for i, h in enumerate(headers):
        cell = tbl.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = INDIGO
        p = cell.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = h
        run.font.name = FONT_MAIN
        run.font.bold = True
        run.font.size = Pt(12)
        run.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    # Data rows
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            if (ri + 1) % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = PALE_INDIGO
            p = cell.text_frame.paragraphs[0]
            run = p.add_run()
            is_ab = val.get('abnormal', False) if isinstance(val, dict) else False
            text  = val['text'] if isinstance(val, dict) else str(val)
            run.text = text
            run.font.name = FONT_MAIN
            run.font.size = Pt(12)
            run.font.color.rgb = RED if is_ab else BLACK
            run.font.bold = is_ab
            p.alignment = PP_ALIGN.CENTER


# ─────────────────────────────────────────────
# Helper: highlighted info box
# ─────────────────────────────────────────────
def add_info_box(slide, left, top, width, height, fill_color=PALE_INDIGO, border_color=INDIGO):
    box = slide.shapes.add_shape(1, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = fill_color
    box.line.color.rgb = border_color
    return box


# ═══════════════════════════════════════════════════════════
# SLIDE 1: Title
# ═══════════════════════════════════════════════════════════
s = add_slide()

# Full INDIGO background
add_bg(s, INDIGO)

# Large decorative light-blue block (right side)
deco = s.shapes.add_shape(1, Inches(9.8), 0, Inches(3.53), prs.slide_height)
deco.fill.solid()
deco.fill.fore_color.rgb = RGBColor(36, 51, 124)   # slightly darker indigo
deco.line.fill.background()

# Thin accent line
acc = s.shapes.add_shape(1, Inches(0.5), Inches(2.85), Inches(9.0), Pt(2))
acc.fill.solid(); acc.fill.fore_color.rgb = LIGHT_BLUE; acc.line.fill.background()

add_textbox(s, "A Case of ST-Elevation Myocardial Infarction",
            Inches(0.5), Inches(1.3), Inches(9.0), Inches(1.2),
            font_size=30, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

add_textbox(s, "STEMI — Anterior Wall  |  LAD Occlusion",
            Inches(0.5), Inches(2.5), Inches(9.0), Inches(0.5),
            font_size=17, bold=False, color=LIGHT_BLUE, align=PP_ALIGN.LEFT)

add_textbox(s, "Department of Cardiology / Internal Medicine",
            Inches(0.5), Inches(3.15), Inches(9.0), Inches(0.45),
            font_size=14, color=WHITE, align=PP_ALIGN.LEFT)

add_textbox(s, "Medical Student Case Presentation  |  April 2026",
            Inches(0.5), Inches(3.65), Inches(9.0), Inches(0.4),
            font_size=13, color=LIGHT_BLUE, align=PP_ALIGN.LEFT)

# Right panel label
add_textbox(s, "CASE\nPRESENTATION",
            Inches(10.1), Inches(2.8), Inches(3.0), Inches(2.0),
            font_size=28, bold=True, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════
# SLIDE 2: Patient Information
# ═══════════════════════════════════════════════════════════
s = add_slide()
add_bg(s)
add_header(s, "01", "Patient Information")

add_table(s,
    ["Field", "Details"],
    [
        ["Age / Sex", "62 y/o Male"],
        ["Chief Complaint", "Sudden onset crushing chest pain — 2 hours"],
        ["Date of Visit", "Emergency admission via 119 ambulance"],
        ["Triage Level", {"text": "ESI Level 1 (Immediate)", "abnormal": True}],
        ["Vital Signs on Arrival", {"text": "BP 158/96  |  HR 102  |  RR 22  |  SpO₂ 94% (RA)", "abnormal": True}],
    ],
    Inches(0.5), Inches(1.1), Inches(12.3),
    col_widths=[3.2, 9.1]
)

# ═══════════════════════════════════════════════════════════
# SLIDE 3: HPI
# ═══════════════════════════════════════════════════════════
s = add_slide()
add_bg(s)
add_header(s, "02", "History of Present Illness (HPI)")

add_bullets(s, [
    {'text': 'Timeline', 'bold': True, 'size': 15, 'color': INDIGO},
    {'text': '1 week prior  — exertional chest discomfort, ignored', 'level': 1, 'size': 13},
    {'text': '2 hours before arrival  — sudden substernal pain at rest, progressive worsening', 'level': 1, 'size': 13},
    {'text': 'Not relieved by position or rest', 'level': 1, 'size': 13},
    {'text': '', 'size': 6},
    {'text': 'OLDCARTS', 'bold': True, 'size': 15, 'color': INDIGO},
], Inches(0.5), Inches(1.1), Inches(5.5), Inches(5.8))

add_table(s,
    ["Feature", "Description"],
    [
        ["Onset",         "Sudden, at rest"],
        ["Location",      "Substernal / precordial"],
        ["Duration",      {"text": "~2 hrs, persistent", "abnormal": True}],
        ["Character",     {"text": "Crushing, pressure-like", "abnormal": True}],
        ["Radiation",     {"text": "Left arm, jaw", "abnormal": True}],
        ["Severity",      {"text": "9/10", "abnormal": True}],
        ["Associated Sx", "Diaphoresis, nausea, dyspnea, lightheadedness"],
    ],
    Inches(6.0), Inches(1.1), Inches(7.0),
    col_widths=[2.2, 4.8]
)

# ═══════════════════════════════════════════════════════════
# SLIDE 4: PMH / FH / SH
# ═══════════════════════════════════════════════════════════
s = add_slide()
add_bg(s)
add_header(s, "03", "Past Medical / Family / Social History")

add_bullets(s, [
    {'text': 'Past Medical History', 'bold': True, 'size': 15, 'color': INDIGO},
    {'text': 'Hypertension (10 yrs) — Amlodipine 5mg QD', 'level': 1, 'size': 13},
    {'text': 'T2DM (7 yrs) — Metformin 1000mg BID', 'level': 1, 'size': 13},
    {'text': 'Hyperlipidemia (5 yrs) — Rosuvastatin 10mg  (poor compliance)', 'level': 1, 'size': 13},
    {'text': 'No prior cardiac Hx, no prior surgeries  |  NKDA', 'level': 1, 'size': 13},
    {'text': '', 'size': 5},
    {'text': 'Family History', 'bold': True, 'size': 15, 'color': INDIGO},
    {'text': 'Father: MI at 55 y/o → deceased', 'level': 1, 'size': 13, 'color': RED},
    {'text': 'Brother: CAD, s/p PCI at 58 y/o', 'level': 1, 'size': 13, 'color': RED},
    {'text': '', 'size': 5},
    {'text': 'Social History', 'bold': True, 'size': 15, 'color': INDIGO},
    {'text': 'Smoking: 30 pack-years (current smoker)', 'level': 1, 'size': 13, 'color': RED},
    {'text': 'Alcohol: 3–4 drinks/week', 'level': 1, 'size': 13},
    {'text': 'Occupation: office worker  |  Lives with wife', 'level': 1, 'size': 13},
], Inches(0.5), Inches(1.1), Inches(12.5), Inches(6.0))

# ═══════════════════════════════════════════════════════════
# SLIDE 5: ROS
# ═══════════════════════════════════════════════════════════
s = add_slide()
add_bg(s)
add_header(s, "04", "Review of Systems (ROS)")

add_table(s,
    ["System", "Positive (+)", "Negative (−)"],
    [
        ["Constitutional", {"text": "Diaphoresis, fatigue (1 wk)", "abnormal": True}, "No fever, no weight loss"],
        ["HEENT",          "—", "Unremarkable"],
        ["Cardiovascular", {"text": "Chest pain, exertional dyspnea", "abnormal": True}, "No palpitations, no syncope"],
        ["Pulmonary",      {"text": "Dyspnea", "abnormal": True}, "No cough, no hemoptysis"],
        ["GI",             {"text": "Nausea", "abnormal": True}, "No vomiting, no abd pain"],
        ["GU",             "—", "Unremarkable"],
        ["MSK",            "—", "Unremarkable"],
        ["Neuro",          {"text": "Lightheadedness", "abnormal": True}, "No syncope, no focal deficit"],
        ["Skin",           "—", "Unremarkable"],
    ],
    Inches(0.5), Inches(1.1), Inches(12.3),
    col_widths=[2.1, 4.4, 5.8]
)

# ═══════════════════════════════════════════════════════════
# SLIDE 6: Physical Exam
# ═══════════════════════════════════════════════════════════
s = add_slide()
add_bg(s)
add_header(s, "05", "Physical Examination")

add_table(s,
    ["BP", "HR", "RR", "Temp", "SpO₂"],
    [[
        {"text": "158/96 mmHg", "abnormal": True},
        {"text": "102 bpm",     "abnormal": True},
        {"text": "22 /min",     "abnormal": True},
        "36.8 °C",
        {"text": "94% (RA)",    "abnormal": True},
    ]],
    Inches(0.5), Inches(1.1), Inches(12.3),
    col_widths=[2.8, 2.2, 2.2, 2.2, 2.9]
)

add_table(s,
    ["System", "Findings"],
    [
        ["General",         {"text": "Acute distress, diaphoretic, pale, clutching chest", "abnormal": True}],
        ["HEENT / Neck",    "No JVD, no carotid bruit"],
        ["Chest / Lungs",   {"text": "Bilateral basal crackles (mild)", "abnormal": True}],
        ["Heart",           {"text": "Tachycardic; S3 gallop; no murmur", "abnormal": True}],
        ["Abdomen",         "Soft, non-tender, no organomegaly"],
        ["Extremities",     "No edema; distal pulses 2+ bilaterally"],
        ["Neuro",           "Alert & oriented ×3; no focal deficits"],
    ],
    Inches(0.5), Inches(2.38), Inches(12.3),
    col_widths=[2.5, 9.8]
)

# ═══════════════════════════════════════════════════════════
# SLIDE 7: Labs & Imaging
# ═══════════════════════════════════════════════════════════
s = add_slide()
add_bg(s)
add_header(s, "06", "Laboratory & Imaging Results")

add_table(s,
    ["Lab", "Result", "Ref", ""],
    [
        ["Troponin I",  {"text": "18.4 ng/mL", "abnormal": True},   "<0.04",   {"text": "↑↑↑", "abnormal": True}],
        ["CK-MB",       {"text": "98 U/L",      "abnormal": True},   "<25",     {"text": "↑↑",  "abnormal": True}],
        ["BNP",         {"text": "420 pg/mL",   "abnormal": True},   "<100",    {"text": "↑",   "abnormal": True}],
        ["WBC",         {"text": "13.2 ×10³/μL","abnormal": True},   "4–10",    "↑"],
        ["Glucose",     {"text": "198 mg/dL",   "abnormal": True},   "70–100",  "↑"],
        ["LDL",         {"text": "168 mg/dL",   "abnormal": True},   "<100",    "↑"],
        ["HDL",         {"text": "32 mg/dL",    "abnormal": True},   ">40",     "↓"],
        ["Creatinine",  "1.1 mg/dL",  "0.7–1.2",  "WNL"],
    ],
    Inches(0.5), Inches(1.1), Inches(6.1),
    col_widths=[2.0, 1.9, 1.3, 0.9]
)

add_bullets(s, [
    {'text': '12-Lead ECG', 'bold': True, 'size': 14, 'color': INDIGO},
    {'text': 'ST elevation 2–4 mm in V1–V4',            'level': 1, 'size': 13, 'color': RED, 'bold': True},
    {'text': 'Reciprocal ST depression in II, III, aVF','level': 1, 'size': 13},
    {'text': 'Q waves developing in V1–V2',             'level': 1, 'size': 13},
    {'text': '→ Anterior STEMI (LAD territory)',        'level': 1, 'size': 13, 'color': RED, 'bold': True},
    {'text': '', 'size': 5},
    {'text': 'Chest X-Ray', 'bold': True, 'size': 14, 'color': INDIGO},
    {'text': 'Mild cardiomegaly',                'level': 1, 'size': 13},
    {'text': 'Pulmonary vascular congestion',    'level': 1, 'size': 13, 'color': RED},
    {'text': 'Early interstitial edema',         'level': 1, 'size': 13, 'color': RED},
    {'text': 'No PTX, no pleural effusion',      'level': 1, 'size': 13},
], Inches(6.5), Inches(1.1), Inches(6.5), Inches(5.8))

# ═══════════════════════════════════════════════════════════
# SLIDE 8: Problem List & Assessment
# ═══════════════════════════════════════════════════════════
s = add_slide()
add_bg(s)
add_header(s, "07", "Problem List & Assessment")

add_bullets(s, [
    {'text': 'Problem List', 'bold': True, 'size': 16, 'color': INDIGO},
    {'text': '1.  Anterior STEMI — LAD occlusion (PRIMARY)', 'level': 1, 'color': RED, 'bold': True, 'size': 14},
    {'text': '2.  Acute heart failure — Killip Class II',    'level': 1, 'color': RED, 'size': 14},
    {'text': '3.  Uncontrolled hypertension',                'level': 1, 'size': 14},
    {'text': '4.  Uncontrolled T2DM  (glucose 198)',         'level': 1, 'size': 14},
    {'text': '5.  Hyperlipidemia  (non-compliant)',          'level': 1, 'size': 14},
    {'text': '6.  Active smoker — 30 pack-years',            'level': 1, 'size': 14},
], Inches(0.5), Inches(1.1), Inches(12.3), Inches(3.0))

# Assessment box
box = add_info_box(s, Inches(0.5), Inches(4.45), Inches(12.3), Inches(2.6))
tf = box.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.18)
tf.margin_top  = Inches(0.12)

p = tf.paragraphs[0]
run = p.add_run()
run.text = "Assessment"
run.font.name  = FONT_MAIN
run.font.bold  = True
run.font.size  = Pt(14)
run.font.color.rgb = INDIGO

p2 = tf.add_paragraph()
run2 = p2.add_run()
run2.text = (
    "62 y/o male with HTN, T2DM, hyperlipidemia, 30 pack-year smoking history, and strong family "
    "history of CAD, presenting with 2-hour crushing substernal chest pain radiating to the left arm "
    "with diaphoresis, ST elevation in V1–V4, and markedly elevated troponin I (18.4 ng/mL), "
    "consistent with anterior STEMI secondary to acute LAD occlusion."
)
run2.font.name  = FONT_MAIN
run2.font.size  = Pt(12)
run2.font.color.rgb = BLACK

# ═══════════════════════════════════════════════════════════
# SLIDE 9: Differential Diagnosis
# ═══════════════════════════════════════════════════════════
s = add_slide()
add_bg(s)
add_header(s, "08", "Differential Diagnosis")

add_table(s,
    ["Diagnosis", "Supporting Evidence", "Against"],
    [
        [{"text": "Anterior STEMI ✓", "abnormal": True},
         "ST↑ V1–V4, Trop↑↑↑, crushing CP, diaphoresis, risk factors",
         "—"],
        ["Aortic Dissection",
         "Severe CP, hypertension",
         "No tearing quality, no pulse differential, normal mediastinum CXR"],
        ["Pulmonary Embolism",
         "Dyspnea, tachycardia, SpO₂↓",
         "No pleuritic pain, no DVT, ECG = ST↑ not S1Q3T3"],
        ["NSTEMI / UA",
         "Chest pain, troponin rise",
         "ST elevation present → rules out NSTEMI by definition"],
        ["Takotsubo CMP",
         "Anterior wall motion abnormality",
         "Male sex, no acute emotional stressor, classic atherosclerotic profile"],
    ],
    Inches(0.5), Inches(1.1), Inches(12.3),
    col_widths=[2.7, 4.6, 5.0]
)

# ═══════════════════════════════════════════════════════════
# SLIDE 10: Final Diagnosis
# ═══════════════════════════════════════════════════════════
s = add_slide()
add_bg(s)
add_header(s, "09", "Final Diagnosis")

add_bullets(s, [
    {'text': 'Anterior ST-Elevation Myocardial Infarction (Anterior STEMI)',
     'bold': True, 'size': 22, 'color': RED},
    {'text': 'Culprit Vessel: LAD (proximal occlusion)',
     'bold': True, 'size': 16, 'color': INDIGO},
    {'text': '', 'size': 7},
    {'text': 'Diagnostic Criteria Met  (ACC/AHA Universal Definition of MI)',
     'bold': True, 'size': 15, 'color': INDIGO},
    {'text': '✅  Ischemic symptoms > 20 minutes',                          'level': 1, 'size': 14},
    {'text': '✅  New ST elevation ≥ 2 mm in V1–V4  (≥2 contiguous leads)', 'level': 1, 'size': 14},
    {'text': '✅  Troponin I 18.4 ng/mL  (> 450× upper limit of normal)',   'level': 1, 'size': 14},
    {'text': '✅  Anterior wall hypokinesia on echocardiogram',              'level': 1, 'size': 14},
    {'text': '', 'size': 7},
    {'text': 'Classification', 'bold': True, 'size': 15, 'color': INDIGO},
    {'text': 'Killip Class II  (S3 gallop + basal crackles)',          'level': 1, 'size': 14},
    {'text': 'TIMI Risk Score: 6/14  →  High Risk',                    'level': 1, 'size': 14, 'color': RED},
    {'text': 'Type 1 MI  (Spontaneous — atherosclerotic plaque rupture)','level': 1, 'size': 14},
], Inches(0.5), Inches(1.1), Inches(12.3), Inches(6.0))

# ═══════════════════════════════════════════════════════════
# SLIDE 11: Treatment Plan
# ═══════════════════════════════════════════════════════════
s = add_slide()
add_bg(s)
add_header(s, "10", "Treatment Plan")

add_bullets(s, [
    {'text': 'Acute / Emergency Management', 'bold': True, 'size': 14, 'color': RED},
    {'text': 'Primary PCI — LAD stent (DES)  |  Door-to-balloon target: < 90 min',
     'level': 1, 'size': 12, 'bold': True},
    {'text': 'O₂ supplementation (target SpO₂ ≥ 95%)  |  IV access ×2  |  Cardiac monitoring  |  Serial ECG',
     'level': 1, 'size': 12},
], Inches(0.5), Inches(1.1), Inches(12.3), Inches(1.4))

add_table(s,
    ["Drug", "Dose / Regimen", "Route", "Indication"],
    [
        ["Aspirin",              "300 mg load → 100 mg QD",       "PO", "Antiplatelet"],
        ["Clopidogrel",          "600 mg load → 75 mg QD",        "PO", "DAPT (×12 mo)"],
        ["Heparin UFH",          "60 U/kg bolus → infusion",      "IV",  "Anticoagulation"],
        ["Metoprolol succinate", "25 mg QD → uptitrate",          "PO", "Beta-blocker"],
        ["Ramipril",             "2.5 mg QD → uptitrate",         "PO", "ACEi (EF↓, HTN)"],
        ["Rosuvastatin",         "20 mg QD",                      "PO", "High-intensity statin"],
    ],
    Inches(0.5), Inches(2.55), Inches(12.3),
    col_widths=[2.6, 3.5, 1.4, 4.8]
)

add_bullets(s, [
    {'text': 'Non-pharmacological / Follow-up', 'bold': True, 'size': 13, 'color': INDIGO},
    {'text': 'Low-Na/low-fat diet  |  Smoking cessation  |  Cardiac rehab referral', 'level': 1, 'size': 12},
    {'text': 'Echo at 4–6 wks  |  Outpatient Cardiology at 2 wks  |  Lipid + HbA1c at 6 wks', 'level': 1, 'size': 12},
], Inches(0.5), Inches(6.25), Inches(12.3), Inches(1.1))

# ═══════════════════════════════════════════════════════════
# SLIDE 12: Discussion & Key Takeaways
# ═══════════════════════════════════════════════════════════
s = add_slide()
add_bg(s)
add_header(s, "11", "Discussion & Key Takeaways")

add_bullets(s, [
    {'text': '"Time is Muscle"', 'bold': True, 'size': 16, 'color': RED},
    {'text': 'Door-to-balloon < 90 min is the #1 determinant of outcome in STEMI', 'level': 1, 'size': 13},
    {'text': 'Every 30-min delay ≈ 7.5% increase in 1-year mortality', 'level': 1, 'size': 13},
    {'text': '', 'size': 5},
    {'text': 'Killip Classification guides acute HF management', 'bold': True, 'size': 16, 'color': INDIGO},
    {'text': 'Class I: No HF  |  Class II: S3/crackles  |  Class III: Pulm edema  |  Class IV: Cardiogenic shock',
     'level': 1, 'size': 12},
    {'text': '', 'size': 5},
    {'text': 'DAPT Duration Post-DES: minimum 12 months', 'bold': True, 'size': 16, 'color': INDIGO},
    {'text': 'Aspirin + P2Y12 inhibitor  (prefer ticagrelor > clopidogrel per guidelines)', 'level': 1, 'size': 12},
    {'text': '', 'size': 5},
    {'text': 'ACEi / ARB mandatory when EF < 40% post-MI', 'bold': True, 'size': 16, 'color': INDIGO},
    {'text': 'Reduces ventricular remodeling and mortality  (SAVE, AIRE trials)', 'level': 1, 'size': 12},
    {'text': '', 'size': 5},
    {'text': 'Secondary Prevention = as important as acute therapy', 'bold': True, 'size': 16, 'color': INDIGO},
    {'text': 'Statin (LDL<70) + BP control + DM mgmt + smoking cessation + cardiac rehab → ↓ 50% recurrence',
     'level': 1, 'size': 12},
], Inches(0.5), Inches(1.1), Inches(12.3), Inches(6.0))

# ─────────────────────────────────────────────
# Save
# ─────────────────────────────────────────────
output_path = "/home/user/endo/case_MI.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
